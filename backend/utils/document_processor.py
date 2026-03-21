"""Document processing utilities."""
from typing import Dict, Any, List
import logging
from pathlib import Path
import re

from backend.utils.ocr_processor import OcrProcessor
from backend.utils.pdf_parser_enhanced import EnhancedPDFParser
from backend.utils.format_converters import convert_to_pdf, can_convert, get_supported_extensions
from backend.utils.table_extractor import extract_text_and_tables

logger = logging.getLogger(__name__)

# Expand supported extensions
SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"} | get_supported_extensions()
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


class DocumentProcessor:
    """Handles document parsing and transformation."""
    
    def __init__(self):
        self.pdf_parser = EnhancedPDFParser()
        self._temp_converted_files = []  # Track temp files for cleanup

    async def parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse a document and extract text content."""
        file_path_obj = Path(file_path)
        file_ext = file_path_obj.suffix.lower()
        
        # Auto-convert to PDF if needed
        original_path = file_path
        if file_ext != '.pdf' and can_convert(file_path):
            logger.info(f"Converting {file_ext} to PDF for enhanced processing...")
            converted_pdf = convert_to_pdf(file_path)
            if converted_pdf:
                file_path = converted_pdf
                file_ext = '.pdf'
                self._temp_converted_files.append(converted_pdf)
            else:
                logger.warning(f"Conversion failed for {file_ext}, falling back to native parser")

        parsers = {
            ".txt": self._parse_text_file,
            ".pdf": self._parse_pdf_file,
            ".docx": self._parse_word_file,
            ".pptx": self._parse_pptx_file,
        }

        parser = parsers.get(file_ext)
        if not parser:
            return {
                "text": "",
                "metadata": {"error": f"Unsupported file type: {file_ext}"}
            }
        return await parser(file_path)

    async def _parse_text_file(self, file_path: str) -> Dict[str, Any]:
        """Parse a text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return {
                "text": text,
                "metadata": {"type": "text", "char_count": len(text)}
            }
        except Exception as e:
            logger.error(f"Error parsing text file: {e}")
            return {"text": "", "metadata": {"error": str(e)}}

    async def _parse_pdf_file(self, file_path: str) -> Dict[str, Any]:
        """Parse a PDF file using enhanced PyMuPDF parser with spatial data and table extraction."""
        try:
            # First, try to extract tables using pdfplumber
            logger.info(f"Extracting text and tables from PDF: {file_path}")
            table_result = extract_text_and_tables(file_path)
            
            # Check if pdfplumber extraction was successful
            if table_result and table_result.get("combined_text"):
                # Use enhanced parser for spatial blocks
                spatial_result = self.pdf_parser.parse_pdf(file_path)
                
                # Merge results
                return {
                    "text": table_result["combined_text"],
                    "blocks": spatial_result.get("blocks", []),
                    "tables": table_result.get("tables", []),
                    "metadata": {
                        "type": "pdf",
                        "is_scanned": spatial_result["metadata"].get("is_scanned", False),
                        "char_count": len(table_result["combined_text"]),
                        "has_spatial_data": True,
                        "table_count": len(table_result.get("tables", [])),
                        "text_block_count": len(table_result.get("text_blocks", []))
                    }
                }
            
            # Fallback to basic PyMuPDF if pdfplumber failed
            result = self.pdf_parser.parse_pdf(file_path)
            
            # Check if it's scanned or very little text extracted
            is_scanned = result["metadata"].get("is_scanned", False)
            
            # If scanned or very little text extracted, use OCR
            if is_scanned or len(result["text"].strip()) < 50:
                logger.info(f"Scanned PDF detected or low text yield. Running OCR...")
                text, ocr_conf = OcrProcessor.process_scanned_pdf(file_path)
                return {
                    "text": text,
                    "blocks": [],  # OCR doesn't preserve blocks yet
                    "tables": [],
                    "metadata": {
                        "type": "pdf",
                        "is_scanned": True,
                        "ocr_confidence": ocr_conf,
                        "char_count": len(text),
                        "has_spatial_data": False
                    }
                }
            
            # Return enhanced result with blocks
            return result
            
        except Exception as e:
            logger.error(f"Error parsing PDF file: {e}")
            return {"text": "", "blocks": [], "tables": [], "metadata": {"error": str(e)}}

    async def _parse_word_file(self, file_path: str) -> Dict[str, Any]:
        """Parse a Word document using python-docx."""
        try:
            from docx import Document

            doc = Document(file_path)
            parts: List[str] = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

            text = "\n\n".join(parts)
            return {
                "text": text,
                "metadata": {
                    "type": "docx",
                    "paragraph_count": len(doc.paragraphs),
                    "table_count": len(doc.tables),
                    "char_count": len(text),
                }
            }
        except Exception as e:
            logger.error(f"Error parsing Word file: {e}")
            return {"text": "", "metadata": {"error": str(e)}}

    async def _parse_pptx_file(self, file_path: str) -> Dict[str, Any]:
        """Parse a PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_text: List[str] = []

            for i, slide in enumerate(prs.slides, 1):
                slide_parts: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_parts.append(para.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                slide_parts.append(row_text)
                if slide_parts:
                    slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_parts))

            text = "\n\n".join(slides_text)
            return {
                "text": text,
                "metadata": {
                    "type": "pptx",
                    "slide_count": len(prs.slides),
                    "char_count": len(text),
                }
            }
        except Exception as e:
            logger.error(f"Error parsing PPTX file: {e}")
            return {"text": "", "metadata": {"error": str(e)}}

    async def chunk_text(
        self,
        text: str,
        chunk_size: int = 1000,
        overlap: int = 200,
        blocks: List[Dict[str, Any]] = None
    ) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: Full text content
            chunk_size: Target chunk size in characters
            overlap: Overlap between chunks
            blocks: Optional list of block dictionaries with bbox data
            
        Returns:
            List of chunk dictionaries (text + optional bbox data)
        """
        # If we have blocks with spatial data, use enhanced chunking
        if blocks:
            return self.pdf_parser.chunk_blocks(blocks, chunk_size, overlap)
        
        # Fallback to simple text chunking
        if not text:
            return []

        sentences = re.split(r'(?<=[.!?])\s+', text)

        chunks = []
        current_chunk: List[str] = []
        current_size = 0

        for sentence in sentences:
            sentence_size = len(sentence)

            if current_size + sentence_size > chunk_size and current_chunk:
                chunks.append({"text": ' '.join(current_chunk)})
                # Keep trailing sentences for overlap
                overlap_chunk: List[str] = []
                overlap_size = 0
                for s in reversed(current_chunk):
                    if overlap_size + len(s) > overlap:
                        break
                    overlap_chunk.insert(0, s)
                    overlap_size += len(s)
                current_chunk = overlap_chunk + [sentence]
                current_size = overlap_size + sentence_size
            else:
                current_chunk.append(sentence)
                current_size += sentence_size

        if current_chunk:
            chunks.append({"text": ' '.join(current_chunk)})

        return chunks
